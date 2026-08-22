"""生成演示用样例文档：虚构「星云科技」公司的 5 种格式企业文档。

- 员工手册.txt
- 差旅报销制度.docx（含表格）
- 信息安全管理制度.pdf
- 产品报价单.xlsx
- 年度战略发布会.pptx

运行: python scripts/make_sample_docs.py [输出目录]
"""

from __future__ import annotations

import sys
from pathlib import Path

TRAVEL_POLICY_LINES = [
    "星云科技（NebulaTech）差旅与费用报销管理制度（2025 版）",
    "",
    "一、总则",
    "本制度适用于公司全体员工因公出差产生的交通、住宿、餐饮及市内交通费用的报销。",
    "",
    "二、交通标准",
    "1. 高铁/动车：行程 800 公里以内限二等座，800 公里以上可乘坐一等座。",
    "2. 飞机：仅限行程超过 1500 公里或经部门负责人特批，限经济舱。",
    "3. 市内交通：优先使用公共交通；打车单次超过 80 元需在报销单中注明事由。",
    "",
    "三、住宿标准",
    "一线城市（北京/上海/广州/深圳）住宿上限为每晚 600 元；",
    "其他城市上限为每晚 400 元；超支部分由个人承担。",
    "",
    "四、餐饮补助",
    "出差期间餐饮补助按天发放：一线城市每天 120 元，其他城市每天 80 元，无需发票。",
    "",
    "五、报销流程",
    "1. 出差结束后 10 个工作日内，在 OA 系统提交《差旅报销单》并上传发票。",
    "2. 部门负责人审批（不超过 3 个工作日），财务部复核后 5 个工作日内打款。",
    "3. 单笔报销金额超过 2 万元或涉及海外出差的，须报分管副总裁特批。",
    "",
    "六、违规处理",
    "虚开发票、重复报销的，一经查实按公司《员工纪律管理办法》给予记过直至解除劳动合同处分。",
]

TRAVEL_TABLE_ROWS = [
    ["城市类别", "住宿上限(元/晚)", "餐饮补助(元/天)"],
    ["一线城市（北上广深）", "600", "120"],
    ["其他城市", "400", "80"],
    ["港澳台及海外", "1200", "200（实报实销）"],
]

SECURITY_POLICY_LINES = [
    "星云科技信息安全管理制度",
    "",
    "第一条 目的：规范公司信息系统与数据的安全管理，防范数据泄露与网络攻击。",
    "",
    "第二条 账号管理：员工须使用企业统一身份认证（SSO）登录各系统，密码长度不得少于 12 位，",
    "且必须包含大小写字母、数字和特殊符号；密码每 90 天强制更换一次。",
    "",
    "第三条 数据分级：公司数据分为公开、内部、机密、绝密四级。机密及以上数据不得通过",
    "个人邮箱、网盘或即时通讯工具外发；确需外发的，须经部门负责人和首席安全官双重审批。",
    "",
    "第四条 终端安全：办公电脑必须安装公司统一防病毒软件并开启磁盘加密；禁止私自安装",
    "未经验证的第三方软件；离开工位须锁屏。",
    "",
    "第五条 应急响应：发生疑似数据泄露时，须在 1 小时内报告信息安全部，",
    "由信息安全部牵头启动应急响应流程（止损 → 取证 → 修复 → 复盘）。",
    "",
    "第六条 违规处罚：违反本制度造成数据泄露的，视情节给予警告、降薪、解除劳动合同处分；",
    "构成犯罪的，移交司法机关处理。",
]

PRICE_ROWS = [
    ["产品型号", "产品名称", "单价(元)", "最低起订量", "备注"],
    ["NX-100", "星云边缘网关", "1899", "10", "含三年质保"],
    ["NX-200", "星云边缘网关 Pro", "3299", "5", "含 5G 模组"],
    ["NX-Cam-A", "AI 巡检摄像头", "2599", "20", "支持边缘推理"],
    ["NX-Cloud", "星云云平台订阅", "1500", "1", "每节点/年"],
    ["NX-Edge-SW", "边缘操作系统授权", "800", "50", "按设备授权"],
]

PPT_SLIDES = [
    ["星云科技 2025 年度战略发布会", "主题：边缘智能，连接万物", "演讲人：CEO 林致远"],
    ["年度回顾", "2024 年营收 8.6 亿元，同比增长 42%", "边缘网关出货量突破 50 万台", "服务企业客户 3200 家"],
    ["战略一：边缘 AI 全栈化", "发布 NX-200 Pro 边缘网关，内置 NPU 算力 26 TOPS", "边缘推理延迟降至 5ms 以内"],
    ["战略二：行业深耕", "聚焦智慧工厂、智慧园区、新能源三大行业", "目标：行业市占率进入前三"],
    ["战略三：生态开放", "开放边缘 OS 与 MCP 工具生态", "与 100 家 ISV 建立联合解决方案"],
    ["2025 年度目标", "营收目标 13 亿元，同比增长 51%", "出货量目标 90 万台", "海外收入占比提升至 20%"],
]

MANUAL_LINES = [
    "星云科技员工手册（节选）",
    "",
    "第一条 工作时间：公司实行弹性工作制，核心工作时间为 10:00-16:00，",
    "员工每日工作时间不少于 8 小时。",
    "",
    "第二条 考勤与请假：年假每年 10 天起，工龄每满 1 年增加 1 天，上限 15 天；",
    "病假需提供二级甲等以上医院证明。",
    "",
    "第三条 薪酬发放：每月 10 日发放上月工资；绩效奖金按季度考核，于次季度首月随工资发放。",
    "",
    "第四条 试用期：新员工试用期 3 个月，表现优秀可提前至 1 个月转正。",
    "",
    "第五条 员工福利：五险一金按国家规定足额缴纳；补充商业医疗保险覆盖员工及一名直系亲属；",
    "每年一次免费体检。",
    "",
    "第六条 离职管理：员工主动离职须提前 30 天书面通知；离职交接包括工作、资产与数据权限回收。",
]


def write_txt(path: Path, lines: list[str]) -> None:
    path.write_text("\n".join(lines), encoding="utf-8")
    print(f"  [TXT]  {path.name}")


def write_docx(path: Path) -> None:
    import docx

    doc = docx.Document()
    doc.add_heading("星云科技差旅与费用报销管理制度", level=1)
    for line in TRAVEL_POLICY_LINES[1:]:
        if line.startswith(("一、", "二、", "三、", "四、", "五、", "六、")):
            doc.add_heading(line, level=2)
        elif line:
            doc.add_paragraph(line)
    doc.add_heading("附表：住宿与餐饮标准", level=2)
    table = doc.add_table(rows=len(TRAVEL_TABLE_ROWS), cols=3)
    table.style = "Table Grid"
    for i, row in enumerate(TRAVEL_TABLE_ROWS):
        for j, cell in enumerate(row):
            table.rows[i].cells[j].text = cell
    doc.save(str(path))
    print(f"  [DOCX] {path.name}")


def write_pdf(path: Path) -> None:
    from fpdf import FPDF

    # CJK 字体：Windows 自带 simhei.ttf，Linux 常见文泉驿/Noto，找不到则报错提示
    candidates = [
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\msyh.ttc"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"),
    ]
    font_path = next((p for p in candidates if p.exists()), None)
    if font_path is None:
        raise FileNotFoundError(
            "未找到 CJK 字体（simhei.ttf / Noto Sans CJK 等），无法生成中文 PDF。"
            "请安装字体或修改 scripts/make_sample_docs.py 的 candidates 列表。"
        )
    pdf = FPDF()
    pdf.add_font("cjk", "", str(font_path))
    pdf.add_page()
    pdf.set_font("cjk", "", 16)
    pdf.cell(0, 10, "星云科技信息安全管理制度", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("cjk", "", 11)
    for line in SECURITY_POLICY_LINES[1:]:
        pdf.multi_cell(0, 6, line)
        pdf.ln(1)
    pdf.output(str(path))
    print(f"  [PDF]  {path.name}")


def write_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    ws = wb.active
    ws.title = "2025年产品报价单"
    for row in PRICE_ROWS:
        ws.append(row)
    ws["A1"].font = Font(bold=True)
    for col in ws.columns:
        width = max(len(str(c.value or "")) for c in col) + 4
        ws.column_dimensions[col[0].column_letter].width = min(width, 40)
    wb.save(str(path))
    print(f"  [XLSX] {path.name}")


def write_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    for slide_lines in PPT_SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title, body = slide.shapes.title, slide.placeholders[1]
        title.text = slide_lines[0]
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(slide_lines[1:]):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(18)
    prs.save(str(path))
    print(f"  [PPTX] {path.name}")


def main() -> None:
    out_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("samples")
    out_dir.mkdir(parents=True, exist_ok=True)
    print(f"生成样例文档到 {out_dir.resolve()}:")
    write_txt(out_dir / "员工手册.txt", MANUAL_LINES)
    write_docx(out_dir / "差旅报销制度.docx")
    write_pdf(out_dir / "信息安全管理制度.pdf")
    write_xlsx(out_dir / "产品报价单.xlsx")
    write_pptx(out_dir / "年度战略发布会.pptx")
    print("完成。")


if __name__ == "__main__":
    main()
