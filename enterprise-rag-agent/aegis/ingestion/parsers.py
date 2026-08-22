"""多格式文档解析：PDF / DOCX / XLSX / PPTX / TXT / MD。

每个解析器把文件拆成若干「原始片段」，附带来源元数据（页码/表格/幻灯片）。
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

SUPPORTED_EXTS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".csv"}


@dataclass
class RawSegment:
    text: str
    metadata: dict = field(default_factory=dict)


class ParseError(RuntimeError):
    pass


def parse_file(path: str | Path) -> list[RawSegment]:
    path = Path(path)
    if not path.exists():
        raise ParseError(f"文件不存在: {path}")
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise ParseError(f"不支持的文件格式: {ext}（支持 {sorted(SUPPORTED_EXTS)}）")
    if path.stat().st_size > 100 * 1024 * 1024:
        raise ParseError("文件超过 100MB 限制")
    if ext == ".pdf":
        return _parse_pdf(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext == ".xlsx":
        return _parse_xlsx(path)
    if ext in (".csv",):
        return _parse_csv(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    return _parse_text(path)


# ---------------- TXT / MD ----------------

def _parse_text(path: Path) -> list[RawSegment]:
    for enc in ("utf-8", "gb18030"):
        try:
            text = path.read_text(encoding=enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        raise ParseError(f"无法识别文本编码: {path.name}")
    text = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        raise ParseError(f"文件内容为空: {path.name}")
    return [RawSegment(text=text, metadata={"source": path.name, "kind": "text"})]


# ---------------- PDF ----------------

def _parse_pdf(path: Path) -> list[RawSegment]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    segs: list[RawSegment] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:  # noqa: BLE001 —— 单页失败不阻断整体
            text = ""
        text = _clean(text)
        if text:
            segs.append(RawSegment(text=text, metadata={"source": path.name, "kind": "pdf", "page": i}))
    if not segs:
        raise ParseError(f"PDF 未提取到文本（可能为扫描件）: {path.name}")
    return segs


# ---------------- DOCX ----------------

def _parse_docx(path: Path) -> list[RawSegment]:
    import docx

    doc = docx.Document(str(path))
    segs: list[RawSegment] = []
    parts: list[str] = []

    def flush(meta: dict) -> None:
        text = _clean("\n".join(parts))
        parts.clear()
        if text:
            segs.append(RawSegment(text=text, metadata=meta))

    # 按文档顺序遍历块级元素（python-docx 不直接支持，用 XML 级联遍历）
    from docx.document import Document as _Doc
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    def iter_block_items(parent: _Doc):
        for child in parent.element.body.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            parts.append(block.text)
        else:  # Table
            flush({"source": path.name, "kind": "docx"})
            rows: list[str] = []
            for row in block.rows:
                cells = [_clean(c.text) for c in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                segs.append(
                    RawSegment(text="\n".join(rows), metadata={"source": path.name, "kind": "docx", "element": "table"})
                )
    flush({"source": path.name, "kind": "docx"})
    if not segs:
        raise ParseError(f"DOCX 未提取到文本: {path.name}")
    return segs


# ---------------- XLSX ----------------

def _parse_xlsx(path: Path) -> list[RawSegment]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    segs: list[RawSegment] = []
    for ws in wb.worksheets:
        lines: list[str] = [f"# 工作表: {ws.title}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
        if len(lines) > 1:
            segs.append(
                RawSegment(
                    text="\n".join(lines),
                    metadata={"source": path.name, "kind": "xlsx", "sheet": ws.title},
                )
            )
    wb.close()
    if not segs:
        raise ParseError(f"XLSX 未提取到内容: {path.name}")
    return segs


def _parse_csv(path: Path) -> list[RawSegment]:
    return _parse_text(path)  # CSV 按文本处理，行结构自然保留


# ---------------- PPTX ----------------

def _parse_pptx(path: Path) -> list[RawSegment]:
    from pptx import Presentation

    prs = Presentation(str(path))
    segs: list[RawSegment] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        lines.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        text = _clean("\n".join(lines))
        if text:
            segs.append(RawSegment(text=text, metadata={"source": path.name, "kind": "pptx", "slide": i}))
    if not segs:
        raise ParseError(f"PPTX 未提取到文本: {path.name}")
    return segs


def _clean(text: str) -> str:
    text = text.replace("\u3000", " ").replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
