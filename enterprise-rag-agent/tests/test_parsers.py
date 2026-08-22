"""多格式解析器测试（真实文件，非 mock）。"""

from pathlib import Path

import pytest

from aegis.ingestion.parsers import parse_file


def _make_docx(path: Path) -> None:
    import docx

    doc = docx.Document()
    doc.add_heading("测试标题", level=1)
    doc.add_paragraph("第一段：差旅住宿标准是每晚 600 元。")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "城市"
    table.rows[0].cells[1].text = "上限"
    table.rows[1].cells[0].text = "北京"
    table.rows[1].cells[1].text = "600"
    doc.save(str(path))


def _make_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["型号", "价格"])
    ws.append(["NX-100", "1899"])
    wb.save(str(path))


def _make_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "战略发布会"
    slide.placeholders[1].text = "营收目标 13 亿元"
    prs.save(str(path))


def _make_pdf(path: Path) -> None:
    from fpdf import FPDF

    candidates = [
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"),
    ]
    font_path = next((p for p in candidates if p.exists()), None)
    assert font_path is not None, "未找到 CJK 字体"
    pdf = FPDF()
    pdf.add_font("cjk", "", str(font_path))
    pdf.add_page()
    pdf.set_font("cjk", size=12)
    pdf.multi_cell(0, 6, "信息安全制度：密码每 90 天更换一次。")
    pdf.output(str(path))


def test_txt(tmp_path: Path):
    p = tmp_path / "a.txt"
    p.write_text("员工手册\n第一条：年假 10 天起。", encoding="utf-8")
    segs = parse_file(p)
    assert len(segs) == 1
    assert "年假 10 天" in segs[0].text
    assert segs[0].metadata["kind"] == "text"


def test_docx(tmp_path: Path):
    p = tmp_path / "a.docx"
    _make_docx(p)
    segs = parse_file(p)
    joined = "\n".join(s.text for s in segs)
    assert "600 元" in joined
    assert "北京" in joined and "上限" in joined  # 表格内容被提取


def test_xlsx(tmp_path: Path):
    p = tmp_path / "a.xlsx"
    _make_xlsx(p)
    segs = parse_file(p)
    joined = "\n".join(s.text for s in segs)
    assert "NX-100" in joined and "1899" in joined
    assert segs[0].metadata["kind"] == "xlsx"


def test_pptx(tmp_path: Path):
    p = tmp_path / "a.pptx"
    _make_pptx(p)
    segs = parse_file(p)
    joined = "\n".join(s.text for s in segs)
    assert "战略发布会" in joined and "13 亿元" in joined


def test_pdf(tmp_path: Path):
    p = tmp_path / "a.pdf"
    _make_pdf(p)
    segs = parse_file(p)
    joined = "\n".join(s.text for s in segs)
    assert "90 天" in joined
    assert segs[0].metadata["page"] == 1


def test_unsupported_extension(tmp_path: Path):
    p = tmp_path / "a.bin"
    p.write_bytes(b"x")
    with pytest.raises(Exception):
        parse_file(p)
